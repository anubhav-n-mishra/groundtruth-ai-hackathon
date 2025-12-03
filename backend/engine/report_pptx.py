"""
PowerPoint report generation module for the Automated Insight Engine.

This module creates downloadable PPTX reports using python-pptx,
including cover slides, charts, tables, QR codes, and recommendations.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from typing import Dict, List, Any, Optional, Tuple
from pathlib import Path
from datetime import datetime
import uuid
import sys
import os
import io
import secrets

# Add parent directory to path for imports
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from core.logger import get_report_logger
from engine.narrative import NarrativeSection

logger = get_report_logger()


# Color scheme
COLORS = {
    "primary": RGBColor(0, 82, 147),      # Dark blue
    "secondary": RGBColor(0, 150, 199),   # Light blue
    "accent": RGBColor(255, 153, 0),      # Orange
    "text_dark": RGBColor(51, 51, 51),    # Dark gray
    "text_light": RGBColor(255, 255, 255), # White
    "background": RGBColor(245, 245, 245), # Light gray
    "success": RGBColor(40, 167, 69),     # Green
    "danger": RGBColor(220, 53, 69),      # Red
}

CHART_COLORS = [
    RGBColor(0, 82, 147),
    RGBColor(0, 150, 199),
    RGBColor(255, 153, 0),
    RGBColor(40, 167, 69),
    RGBColor(148, 103, 189),
    RGBColor(220, 53, 69),
]


def create_presentation() -> Presentation:
    """Create a new PowerPoint presentation with default settings."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 widescreen
    prs.slide_height = Inches(7.5)
    return prs


def add_header_bar(slide, prs, title: str, color=None):
    """Add a header bar with title to a slide."""
    if color is None:
        color = COLORS["primary"]
    
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(1.0)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = color
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.25),
        Inches(12), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS["text_light"]


def add_cover_slide(
    prs: Presentation,
    title: str = "Automated Insight Engine",
    subtitle: Optional[str] = None,
    date_str: Optional[str] = None
) -> None:
    """Add a cover slide to the presentation."""
    logger.debug("Adding cover slide")
    
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS["primary"]
    background.line.fill.background()
    
    # Accent bar at top
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(0.15)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLORS["accent"]
    accent_bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.2),
        Inches(12.333), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS["text_light"]
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3.9),
            Inches(12.333), Inches(0.8)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = subtitle
        subtitle_para.font.size = Pt(22)
        subtitle_para.font.color.rgb = COLORS["secondary"]
        subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Date
    if date_str is None:
        date_str = datetime.now().strftime("%B %d, %Y")
    
    date_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(6.3),
        Inches(12.333), Inches(0.5)
    )
    date_frame = date_box.text_frame
    date_para = date_frame.paragraphs[0]
    date_para.text = f"Generated: {date_str}"
    date_para.font.size = Pt(14)
    date_para.font.color.rgb = COLORS["text_light"]
    date_para.alignment = PP_ALIGN.CENTER
    
    # Powered by label
    powered_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(6.8),
        Inches(12.333), Inches(0.4)
    )
    powered_frame = powered_box.text_frame
    powered_para = powered_frame.paragraphs[0]
    powered_para.text = "Powered by Gemini AI"
    powered_para.font.size = Pt(12)
    powered_para.font.italic = True
    powered_para.font.color.rgb = COLORS["secondary"]
    powered_para.alignment = PP_ALIGN.CENTER


def add_executive_summary_slide(
    prs: Presentation,
    headline: str,
    key_stats: List[Tuple[str, str, str]] = None
) -> None:
    """Add an executive summary slide with KPI cards."""
    logger.debug("Adding executive summary slide")
    
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    add_header_bar(slide, prs, "Executive Summary")
    
    # Headline box
    headline_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.3),
        Inches(12.333), Inches(1.5)
    )
    headline_box.fill.solid()
    headline_box.fill.fore_color.rgb = COLORS["background"]
    headline_box.line.color.rgb = COLORS["secondary"]
    headline_box.line.width = Pt(2)
    
    headline_text = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.6),
        Inches(11.733), Inches(1.0)
    )
    headline_frame = headline_text.text_frame
    headline_frame.word_wrap = True
    headline_para = headline_frame.paragraphs[0]
    headline_para.text = headline
    headline_para.font.size = Pt(22)
    headline_para.font.bold = True
    headline_para.font.color.rgb = COLORS["text_dark"]
    headline_para.alignment = PP_ALIGN.CENTER
    
    # KPI Cards
    if key_stats and len(key_stats) > 0:
        num_cards = min(len(key_stats), 4)
        card_width = 2.8
        total_width = num_cards * card_width + (num_cards - 1) * 0.3
        start_x = (13.333 - total_width) / 2
        
        for i, (label, value, trend) in enumerate(key_stats[:4]):
            x = start_x + i * (card_width + 0.3)
            
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(3.2),
                Inches(card_width), Inches(2.2)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = COLORS["text_light"]
            card.line.color.rgb = COLORS["primary"]
            card.line.width = Pt(2)
            
            value_box = slide.shapes.add_textbox(
                Inches(x + 0.1), Inches(3.5),
                Inches(card_width - 0.2), Inches(0.8)
            )
            value_frame = value_box.text_frame
            value_para = value_frame.paragraphs[0]
            value_para.text = str(value)
            value_para.font.size = Pt(32)
            value_para.font.bold = True
            value_para.font.color.rgb = COLORS["primary"]
            value_para.alignment = PP_ALIGN.CENTER
            
            trend_box = slide.shapes.add_textbox(
                Inches(x + 0.1), Inches(4.3),
                Inches(card_width - 0.2), Inches(0.4)
            )
            trend_frame = trend_box.text_frame
            trend_para = trend_frame.paragraphs[0]
            trend_para.text = trend
            trend_para.font.size = Pt(14)
            if "+" in trend:
                trend_para.font.color.rgb = COLORS["success"]
            elif "-" in trend:
                trend_para.font.color.rgb = COLORS["danger"]
            else:
                trend_para.font.color.rgb = COLORS["text_dark"]
            trend_para.alignment = PP_ALIGN.CENTER
            
            label_box = slide.shapes.add_textbox(
                Inches(x + 0.1), Inches(4.8),
                Inches(card_width - 0.2), Inches(0.5)
            )
            label_frame = label_box.text_frame
            label_para = label_frame.paragraphs[0]
            label_para.text = label
            label_para.font.size = Pt(12)
            label_para.font.color.rgb = COLORS["text_dark"]
            label_para.alignment = PP_ALIGN.CENTER


def add_metrics_chart_slide(
    prs: Presentation,
    metrics_data: Dict[str, Any],
    title: str = "Key Metrics Comparison"
) -> None:
    """Add a slide with a bar chart comparing metrics."""
    logger.debug("Adding metrics chart slide")
    
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    add_header_bar(slide, prs, title)
    
    current_metrics = metrics_data.get("current_totals", {})
    previous_metrics = metrics_data.get("previous_totals", {})
    
    if not current_metrics and not previous_metrics:
        msg_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9), Inches(1.5))
        msg_frame = msg_box.text_frame
        msg_para = msg_frame.paragraphs[0]
        msg_para.text = "Insufficient data for chart visualization.\nEnsure your date range contains data."
        msg_para.font.size = Pt(20)
        msg_para.font.color.rgb = COLORS["text_dark"]
        msg_para.alignment = PP_ALIGN.CENTER
        return
    
    all_metrics = set(current_metrics.keys()) | set(previous_metrics.keys())
    metric_totals = [(m, (current_metrics.get(m, 0) or 0) + (previous_metrics.get(m, 0) or 0)) for m in all_metrics]
    metric_totals.sort(key=lambda x: x[1], reverse=True)
    top_metrics = [m[0] for m in metric_totals[:6]]
    
    chart_data = CategoryChartData()
    chart_data.categories = top_metrics
    chart_data.add_series('Current', [current_metrics.get(m, 0) or 0 for m in top_metrics])
    chart_data.add_series('Previous', [previous_metrics.get(m, 0) or 0 for m in top_metrics])
    
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1.4), Inches(12.333), Inches(5.8), chart_data
    ).chart
    
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    
    plot = chart.plots[0]
    if len(plot.series) > 0:
        plot.series[0].format.fill.solid()
        plot.series[0].format.fill.fore_color.rgb = COLORS["primary"]
    if len(plot.series) > 1:
        plot.series[1].format.fill.solid()
        plot.series[1].format.fill.fore_color.rgb = COLORS["secondary"]


def add_dimension_breakdown_chart(
    prs: Presentation,
    insights_data: Dict[str, Any],
    title: str = "Performance by Dimension"
) -> None:
    """Add a pie chart showing breakdown by dimension."""
    logger.debug("Adding dimension breakdown chart slide")
    
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    add_header_bar(slide, prs, title)
    
    insights = insights_data.get("insights", [])
    
    if not insights:
        msg_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9), Inches(1.5))
        msg_frame = msg_box.text_frame
        msg_para = msg_frame.paragraphs[0]
        msg_para.text = "No dimension breakdown available."
        msg_para.font.size = Pt(20)
        msg_para.font.color.rgb = COLORS["text_dark"]
        msg_para.alignment = PP_ALIGN.CENTER
        return
    
    dim_data = {}
    for insight in insights[:10]:
        dims = insight.get("dimensions", {})
        if dims:
            first_dim = list(dims.values())[0]
            current_val = insight.get("current_value", 0) or 0
            dim_data[str(first_dim)] = dim_data.get(str(first_dim), 0) + abs(current_val)
    
    if not dim_data:
        for insight in insights[:6]:
            label = insight.get("metric", "Unknown")
            impact = abs(insight.get("impact_score", 0))
            dim_data[label] = impact
    
    sorted_dims = sorted(dim_data.items(), key=lambda x: x[1], reverse=True)[:6]
    if not sorted_dims:
        return
    
    chart_data = CategoryChartData()
    chart_data.categories = [d[0] for d in sorted_dims]
    chart_data.add_series('Values', [d[1] for d in sorted_dims])
    
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, Inches(2), Inches(1.4), Inches(9), Inches(5.8), chart_data
    ).chart
    
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    
    plot = chart.plots[0]
    for i, point in enumerate(plot.series[0].points):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = CHART_COLORS[i % len(CHART_COLORS)]


def add_insights_table_slide(
    prs: Presentation,
    insights: List[Dict[str, Any]],
    title: str = "Top Insights by Impact"
) -> None:
    """Add a slide with a detailed insights table."""
    if not insights:
        return
    
    logger.debug("Adding insights table slide")
    
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    add_header_bar(slide, prs, title)
    
    num_rows = min(len(insights), 6) + 1
    table = slide.shapes.add_table(num_rows, 4, Inches(0.4), Inches(1.4), Inches(12.5), Inches(5.5)).table
    
    table.columns[0].width = Inches(4)
    table.columns[1].width = Inches(2.5)
    table.columns[2].width = Inches(3)
    table.columns[3].width = Inches(3)
    
    headers = ["Segment", "Metric", "Change", "Impact"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS["primary"]
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(14)
        para.font.bold = True
        para.font.color.rgb = COLORS["text_light"]
        para.alignment = PP_ALIGN.CENTER
    
    for row_idx, insight in enumerate(insights[:6]):
        dims = insight.get("dimensions", {})
        dim_str = ", ".join(f"{k}: {v}" for k, v in dims.items()) if dims else "Overall"
        metric = insight.get("metric", "N/A")
        delta_pct = insight.get("delta_pct", 0)
        direction = insight.get("direction", "flat")
        impact = insight.get("impact_score", 0)
        
        arrow = "↑" if direction == "up" else "↓" if direction == "down" else "→"
        change_str = f"{arrow} {delta_pct:+.1f}%"
        impact_str = f"{'█' * min(int(impact * 5), 10)} {impact:.2f}"
        
        row_data = [dim_str, metric.upper(), change_str, impact_str]
        for col_idx, value in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = value
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS["background"]
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(11)
            if col_idx == 2:
                para.font.color.rgb = COLORS["success"] if direction == "up" else COLORS["danger"] if direction == "down" else COLORS["text_dark"]
                para.font.bold = True


def add_bullet_slide(prs: Presentation, bullets: List[str], title: str = "Detailed Insights") -> None:
    """Add a slide with bullet points."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    add_header_bar(slide, prs, title)
    
    bullet_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.5))
    bullet_frame = bullet_box.text_frame
    bullet_frame.word_wrap = True
    
    for i, text in enumerate(bullets):
        para = bullet_frame.paragraphs[0] if i == 0 else bullet_frame.add_paragraph()
        para.text = f"• {text}"
        para.font.size = Pt(18)
        para.font.color.rgb = COLORS["text_dark"]
        para.space_before = Pt(14)
        para.space_after = Pt(14)


def add_recommendation_slide(prs: Presentation, recommendation: str, title: str = "Recommendation") -> None:
    """Add a recommendation slide."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    add_header_bar(slide, prs, title, COLORS["accent"])
    
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2), Inches(11.333), Inches(3))
    box.fill.solid()
    box.fill.fore_color.rgb = COLORS["background"]
    box.line.color.rgb = COLORS["accent"]
    box.line.width = Pt(4)
    
    icon = slide.shapes.add_textbox(Inches(5.5), Inches(2.3), Inches(2), Inches(0.8))
    icon.text_frame.paragraphs[0].text = "💡"
    icon.text_frame.paragraphs[0].font.size = Pt(40)
    icon.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    rec = slide.shapes.add_textbox(Inches(1.5), Inches(3.2), Inches(10.333), Inches(1.5))
    rec.text_frame.word_wrap = True
    rec.text_frame.paragraphs[0].text = recommendation
    rec.text_frame.paragraphs[0].font.size = Pt(22)
    rec.text_frame.paragraphs[0].font.color.rgb = COLORS["text_dark"]
    rec.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def add_qr_code_slide(
    prs: Presentation, 
    session_id: str, 
    base_url: str = "http://localhost:8000",
    qr_image_path: Optional[str] = None,
    dashboard_url: Optional[str] = None
) -> str:
    """Add a slide with QR code for dashboard access."""
    logger.debug("Adding QR code slide")
    
    from engine.qrcode_gen import generate_qr_for_dashboard
    
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    add_header_bar(slide, prs, "🎯 Access Live Dashboard", COLORS["primary"])
    
    # Use pre-generated QR code if available
    if qr_image_path and Path(qr_image_path).exists():
        qr_path = qr_image_path
        # Use provided dashboard_url or construct a basic one
        if not dashboard_url:
            dashboard_url = f"{base_url}/dashboard/{session_id}"
    else:
        # Generate new QR code with new token
        token = secrets.token_urlsafe(32)
        qr_bytes, dashboard_url = generate_qr_for_dashboard(
            base_url=base_url, 
            session_id=session_id, 
            token=token, 
            size=400
        )
        
        qr_path = Path("tmp") / f"qr_{session_id}.png"
        qr_path.parent.mkdir(parents=True, exist_ok=True)
        with open(qr_path, 'wb') as f:
            f.write(qr_bytes)
        qr_path = str(qr_path)
    
    # Add QR code image
    slide.shapes.add_picture(str(qr_path), Inches(4.7), Inches(1.5), Inches(4), Inches(4))
    
    # Instruction text
    instr = slide.shapes.add_textbox(Inches(1), Inches(5.7), Inches(11.333), Inches(0.6))
    instr.text_frame.paragraphs[0].text = "📱 Scan this QR code to access the live interactive dashboard"
    instr.text_frame.paragraphs[0].font.size = Pt(18)
    instr.text_frame.paragraphs[0].font.color.rgb = COLORS["text_dark"]
    instr.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Features text
    features = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(11.333), Inches(0.5))
    features.text_frame.paragraphs[0].text = "🎙️ Voice Summary  •  📊 Interactive Charts  •  📈 Real-time Insights"
    features.text_frame.paragraphs[0].font.size = Pt(14)
    features.text_frame.paragraphs[0].font.color.rgb = COLORS["secondary"]
    features.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # URL text
    url = slide.shapes.add_textbox(Inches(1), Inches(6.7), Inches(11.333), Inches(0.4))
    url_display = dashboard_url if len(dashboard_url) < 60 else dashboard_url[:60] + "..."
    url.text_frame.paragraphs[0].text = f"Or visit: {url_display}"
    url.text_frame.paragraphs[0].font.size = Pt(10)
    url.text_frame.paragraphs[0].font.color.rgb = COLORS["text_muted"] if "text_muted" in COLORS else RGBColor(128, 128, 128)
    url.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return dashboard_url


def build_report(
    narrative: NarrativeSection,
    insights_data: Dict[str, Any],
    output_dir: str,
    filename: Optional[str] = None,
    include_qr: bool = True,
    base_url: str = "http://localhost:8000",
    session_id: Optional[str] = None,
    qr_code_path: Optional[str] = None
) -> str:
    """Build a complete PowerPoint report with charts and QR code."""
    logger.info("Building PowerPoint report")
    
    prs = create_presentation()
    config = insights_data.get("config", {})
    current_period = config.get("current_period", "")
    previous_period = config.get("previous_period", "")
    subtitle = f"Comparing {current_period} vs {previous_period}" if current_period else None
    
    # Generate session ID if not provided
    if not session_id:
        session_id = str(uuid.uuid4())[:12]
    
    # 1. Cover
    add_cover_slide(prs, title=narrative.title, subtitle=subtitle)
    
    # 2. Executive Summary
    insights = insights_data.get("insights", [])
    key_stats = []
    for insight in insights[:4]:
        metric = insight.get("metric", "Metric")
        current = insight.get("current_value", 0) or 0
        delta = insight.get("delta_pct", 0)
        direction = insight.get("direction", "flat")
        trend = f"↑ +{delta:.1f}%" if direction == "up" else f"↓ {delta:.1f}%" if direction == "down" else f"→ {delta:.1f}%"
        value_str = f"{current/1000000:.1f}M" if current >= 1000000 else f"{current/1000:.1f}K" if current >= 1000 else f"{current:.0f}"
        key_stats.append((metric, value_str, trend))
    add_executive_summary_slide(prs, headline=narrative.headline, key_stats=key_stats)
    
    # 3-4. Charts
    add_metrics_chart_slide(prs, insights_data)
    add_dimension_breakdown_chart(prs, insights_data)
    
    # 5. Table
    if insights:
        add_insights_table_slide(prs, insights)
    
    # 6. Bullets
    add_bullet_slide(prs, bullets=narrative.bullets)
    
    # 7. Recommendation
    add_recommendation_slide(prs, recommendation=narrative.recommendation)
    
    # 8. QR Code
    if include_qr:
        try:
            add_qr_code_slide(prs, session_id=session_id, base_url=base_url, qr_image_path=qr_code_path)
        except Exception as e:
            logger.warning(f"Failed to add QR code slide: {e}")
    
    # Save
    output_path = Path(output_dir)
    output_path.mkdir(parents=True, exist_ok=True)
    if filename is None:
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        unique_id = str(uuid.uuid4())[:8]
        filename = f"insight_report_{timestamp}_{unique_id}.pptx"
    
    full_path = output_path / filename
    prs.save(str(full_path))
    logger.info(f"Report saved to: {full_path}")
    return str(full_path)


def generate_report(
    narrative: NarrativeSection,
    insights_data: Dict[str, Any],
    output_dir: str = "static/reports",
    dashboard_url: Optional[str] = None,
    qr_code_path: Optional[str] = None
) -> str:
    """
    Main entry point for report generation.
    
    Args:
        narrative: NarrativeSection with content
        insights_data: Dictionary containing insights and config
        output_dir: Directory to save reports
        dashboard_url: Optional URL for dashboard access
        qr_code_path: Optional path to pre-generated QR code image
        
    Returns:
        Path to the generated report file
    """
    # Extract session ID from dashboard URL if present
    session_id = None
    if dashboard_url:
        # URL format: http://localhost:8000/dashboard/SESSION_ID?token=...
        import re
        match = re.search(r'/dashboard/([^?]+)', dashboard_url)
        if match:
            session_id = match.group(1)
    
    # Always include QR code - generate session ID if not present
    include_qr = True
    if not session_id:
        session_id = str(uuid.uuid4())[:12]
    
    return build_report(
        narrative, 
        insights_data, 
        output_dir,
        include_qr=include_qr,
        session_id=session_id,
        qr_code_path=qr_code_path
    )

